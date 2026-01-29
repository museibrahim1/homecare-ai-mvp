#!/usr/bin/env python3
"""
Export pitch deck to a PowerPoint (.pptx) without needing a browser.

Usage:
  python3 pitch_deck/export_pptx.py

Output:
  pitch_deck/Homecare-AI-Pitch-Deck.pptx
"""

from __future__ import annotations

from pathlib import Path


def _fit_image_box(img_w_in: float, img_h_in: float, box_w_in: float, box_h_in: float) -> tuple[float, float]:
    """Return (w, h) to fit image into box while preserving aspect."""
    if img_w_in <= 0 or img_h_in <= 0:
        return box_w_in, box_h_in
    img_ratio = img_w_in / img_h_in
    box_ratio = box_w_in / box_h_in
    if img_ratio >= box_ratio:
        # fit by width
        w = box_w_in
        h = box_w_in / img_ratio
    else:
        # fit by height
        h = box_h_in
        w = box_h_in * img_ratio
    return w, h


def main() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    try:
        from PIL import Image
    except Exception:  # pragma: no cover
        Image = None  # type: ignore

    root = Path(__file__).resolve().parent
    screenshots = root / "screenshots"
    out_path = root / "Homecare-AI-Pitch-Deck.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    EMU_PER_INCH = 914400

    # Layouts (varies by template); fall back safely.
    title_slide_layout = prs.slide_layouts[0]
    title_and_content_layout = prs.slide_layouts[1]
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    def add_title_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    def add_bullets_slide(title: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(title_and_content_layout)
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
            p.level = 0

    def add_screenshot_slide(title: str, filename: str, caption: str | None = None) -> None:
        slide = prs.slides.add_slide(blank_layout)
        sw = int(prs.slide_width)
        sh = int(prs.slide_height)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), sw - Inches(1.2), Inches(0.7))
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(34)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        # Optional caption
        y_after_title = 1.15
        if caption:
            cap_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.05), sw - Inches(1.2), Inches(0.4))
            ctf = cap_box.text_frame
            ctf.text = caption
            ctf.paragraphs[0].font.size = Pt(16)
            ctf.paragraphs[0].alignment = PP_ALIGN.LEFT
            y_after_title = 1.55

        img_path = screenshots / filename
        if not img_path.exists():
            # Leave a visible note if the image is missing.
            warn = slide.shapes.add_textbox(Inches(0.6), Inches(y_after_title), sw - Inches(1.2), Inches(1.0))
            wtf = warn.text_frame
            wtf.text = f"Missing screenshot: {img_path.name}"
            wtf.paragraphs[0].font.size = Pt(18)
            return

        # Image placement box
        margin_x = 0.6
        margin_bottom = 0.45
        box_x = int(Inches(margin_x))
        box_y = int(Inches(y_after_title))
        box_w = sw - int(Inches(margin_x * 2))
        box_h = sh - int(Inches(y_after_title)) - int(Inches(margin_bottom))

        if Image:
            with Image.open(img_path) as im:
                iw, ih = im.size
            # Ratio only (pixels are fine)
            img_w = float(iw)
            img_h = float(ih)
            box_w_in = box_w / EMU_PER_INCH
            box_h_in = box_h / EMU_PER_INCH
            fitted_w_in, fitted_h_in = _fit_image_box(img_w, img_h, box_w_in, box_h_in)
            w = int(Inches(fitted_w_in))
            h = int(Inches(fitted_h_in))
        else:
            w = int(box_w)
            h = int(box_h)

        x = box_x + int((box_w - w) / 2)
        y = box_y + int((box_h - h) / 2)
        slide.shapes.add_picture(str(img_path), x, y, width=w, height=h)

    # ---- Slides (mirrors pitch_deck/deck.md) ----
    add_title_slide(
        "Homecare AI",
        "Care assessments in, proposal-ready contracts out.\nAI-powered workflow for home healthcare agencies.",
    )

    add_bullets_slide(
        "The problem",
        [
            "Slow: intake calls → manual notes → manual pricing → manual contract drafting",
            "Inconsistent: variability across coordinators/branches leads to pricing and scope errors",
            "High-stakes: delays and mistakes cost revenue, compliance risk, and client trust",
        ],
    )

    add_bullets_slide(
        "The solution",
        [
            "Transcript ingestion (upload audio or import text)",
            "Services & billables extraction",
            "Contract drafting from templates",
            "Review + edit + export (PDF/CSV) in one place",
        ],
    )

    add_screenshot_slide("Demo login (admin UI)", "01-login.png")
    add_screenshot_slide("Assessments inbox (the work queue)", "02-assessments.png", "Start a new assessment or run a guided demo flow.")
    add_bullets_slide(
        "Workflow: assessment → pipeline → proposal",
        [
            "Start or import an assessment",
            "Run pipeline steps (transcribe → bill → contract)",
            "Review outputs in the right-side panel",
            "Export PDF/CSV for client + billing ops",
        ],
    )
    add_screenshot_slide("Intake capture + pipeline controls", "03-visit-detail.png")
    add_screenshot_slide("Contract generation (human-in-the-loop)", "04-contract-preview.png", "Edit/regenerate, then export.")
    add_screenshot_slide("Client CRM (lightweight, agency-friendly)", "05-clients.png")
    add_screenshot_slide("Reporting & exports", "06-reports.png")
    add_bullets_slide(
        "How it works (high level)",
        [
            "Frontend: Next.js admin UI",
            "Backend: FastAPI API + auth",
            "Pipeline: worker jobs (transcription, extraction, contract generation)",
            "Storage: Postgres (entities) + S3/MinIO (audio)",
        ],
    )
    add_bullets_slide(
        "Business model (initial)",
        [
            "Subscription per agency / location",
            "Tiered plans by usage (visits/month, storage, seats)",
            "Optional onboarding + integrations",
        ],
    )
    add_bullets_slide(
        "What we’re building next",
        [
            "Deeper template controls + clause library",
            "More robust validation + audit trails",
            "Integrations (EMR/CRM, billing systems)",
            "Multi-location analytics and admin tooling",
        ],
    )

    prs.save(out_path)
    print(f"✅ Wrote {out_path}")


if __name__ == "__main__":
    main()

